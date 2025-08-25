import os
import re
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

from llms.aiforce_util import AiforceUtil


class PPTUtil:
    template_info = {}
    have_Image_format = ['1_Title and Content',
                         '2_Title and Content',
                         '3_Title and Content',
                         '1_Title Only_Black',
                         'Title w/Image',
                         'Title w/Image_Black',
                         'Content w/ Product',
                         'Content w/ Product_Black']
    # 解析AI生成的PPT制作方案，结构化组织数据
    @staticmethod
    def get_ppt_info(ppt_txt: str):
        # 根据“--”进行分块
        slides = re.split(r'^---+\s*$', ppt_txt, flags=re.MULTILINE)

        parsed_slides = []
        # 解析块内容
        for slide in slides:
            slide = slide.strip()
            if not slide:
                continue

            # 提取一级标题（幻灯片主标题）
            main_title_match = re.search(r'^#\s+(.*)$', slide, re.MULTILINE)
            main_title = main_title_match.group(1).strip() if main_title_match else "未命名幻灯片"

            # 提取二级标题及其内容
            subsections = []
            # 匹配 ## 标题及其内容（直到下一个 ## 或结束）
            subsection_pattern = r'^##\s+(.*?)(?=^##\s+|\Z)'
            for match in re.finditer(subsection_pattern, slide, re.DOTALL | re.MULTILINE):
                subsection_text = match.group(0)
                subtitle = re.search(r'^##\s+(.*)$', subsection_text, re.MULTILINE).group(1).strip()
                # 提取嵌套列表内容
                list_items,haveImage = PPTUtil.parse_nested_list(subsection_text)
                # 处理包含图片但是版式不正确的情况
                if haveImage:
                    format = PPTUtil.get_format_or_placeholder(subtitle)
                    if format not in PPTUtil.have_Image_format:
                        subtitle = PPTUtil.get_no_format_or_placeholder(subtitle) + "(Title w/Image)"
                subsections.append({
                    'subtitle': subtitle,
                    'list_items': list_items
                })

            parsed_slides.append({
                'main_title': main_title,
                'subsections': subsections
            })

        return parsed_slides
    # #解析嵌套的Markdown无序列表
    @staticmethod
    def parse_nested_list(text):
        list_pattern = r'^([ \t]*)([-\*])\s+(.*)$'
        items = []
        stack = [items]  # 用于跟踪嵌套层级的栈
        haveImage = False
        for line in text.split('\n'):
            # 保留前导空格，去掉尾部空格
            line = line.rstrip()
            list_match = re.match(list_pattern, line)
            if list_match:
                indent, marker, content = list_match.groups()
                indent_level = len(indent.replace('\t', '    '))  # 将制表符转换为4个空格计算缩进
                if "Picture Placeholder" in content:
                    indent_level = 0
                    haveImage = True
                # 调整栈的大小以匹配当前缩进级别
                while len(stack) > indent_level // 2 + 1:
                    stack.pop()
                while len(stack) < indent_level // 2 + 1:
                    new_level = []
                    #处理嵌入列表无父列表的情况
                    if len(stack[-1]) != 0:
                        stack[-1][-1]['children'] = new_level
                        stack.append(new_level)
                    else:
                        # 添加新列表项
                        item = {
                            'content': content.strip(),
                            'children': []
                        }
                        stack[-1].append(item)

                # 添加新列表项
                item = {
                    'content': content.strip(),
                    'children': []
                }
                stack[-1].append(item)

        return items,haveImage
    # 根据结构化数据制作PPT
    @staticmethod
    def get_ppt_file(ppt_info: list,template=0):
        # 获取ppt文件，其中template指定使用的模板
        fileName = "./utils/template/Template"+str(template)+".pptx"
        prs = Presentation(fileName)
        # 设置首页
        PPTUtil.set_home_page(prs=prs)
        for main_info in ppt_info:
            main_title = main_info['main_title']
            subsections = main_info['subsections']
            for sub_info in subsections:
                # 设置章节页
                sub_title = sub_info['subtitle']
                list_items = sub_info['list_items']
                real_sub_title = PPTUtil.get_no_format_or_placeholder(sub_title)
                PPTUtil.set_sub_title_page(main_title,real_sub_title,prs=prs)
                format = PPTUtil.get_format_or_placeholder(sub_title)
                # 设置内容页
                PPTUtil.set_common_page(format,real_sub_title,list_items,prs)
        PPTUtil.set_thank_page(prs=prs)
        return prs
    # 将多级嵌套列表转换成字符串
    @staticmethod
    def get_list_str(result,content,children):
        result+=content
        if len(children)!=0:
            for i in range(len(children)):
                result = PPTUtil.get_list_str(result,"\n  "+children[i]['content'],children[i]['children'])
        return result
    # 设置PPT首页
    @staticmethod
    def set_home_page(title="Presentation title",info="Speaker name|Date",format="Title Slide_White",prs=None):
        # 查找首页母版页
        for layout in prs.slide_layouts:
            if layout.name==format:
                slide = prs.slides.add_slide(layout)
                # 设置首页内容
                for shape in slide.placeholders:
                    if shape.name=='Title 1' and shape.has_text_frame:
                        shape.text = title
                    if shape.name=='Subtitle 2' and shape.has_text_frame:
                        shape.text = info
    # 设置PPT副标题页
    @staticmethod
    def set_sub_title_page(title="Section Header",subhead="Subhead",format="Section Header_White",prs=None):
        # 查找副标题母版页
        for layout in prs.slide_layouts:
            if layout.name==format:
                slide = prs.slides.add_slide(layout)
                # 设置副标题内容
                for shape in slide.placeholders:
                    if shape.name=='Title 1' and shape.has_text_frame:
                        shape.text = title
                    if shape.name=='Text Placeholder 2' and shape.has_text_frame:
                        shape.text = subhead
    # 设置PPT感谢页
    @staticmethod
    def set_thank_page(format="Closing Slide_Black",prs=None):
        # 查找感谢页母版
        for layout in prs.slide_layouts:
            if layout.name==format:
                slide = prs.slides.add_slide(layout)
    # 设置PPT普通内容页
    @staticmethod
    def set_common_page(format="",sub_title="",items=None,prs=None):
        slide = None
        baseLayout = None
        # 查找匹配的母版页
        for layout in prs.slide_layouts:
            if layout.name == format:
                slide = prs.slides.add_slide(layout)
            if layout.name == "标题和内容":
                baseLayout = layout
        # 处理副标题版式标注错误的情况
        if slide is None:
            slide = prs.slides.add_slide(baseLayout)
        # 设置标题
        for shape in slide.placeholders:
            if shape.name == 'Title 1' and shape.has_text_frame:
                shape.text = sub_title
        # 处理内容
        for item in items:
            placeholder = PPTUtil.get_format_or_placeholder(item['content'])
            # 处理占位符没有标注的情况
            if placeholder == "":
                placeholder = 'Content Placeholder 2'
            # 处理序号错误的情况
            miss_flag = True
            for shape in slide.placeholders:
                if shape.name == placeholder:
                    miss_flag = False
            # 填充内容
            for shape in slide.placeholders:
                if shape.name == placeholder or (miss_flag and placeholder[:-2] in shape.name):
                    miss_flag = False
                    # 处理文本内容
                    if "Content Placeholder" in placeholder or "Text Placeholder" in placeholder:
                        item['content'] = PPTUtil.get_no_format_or_placeholder(item['content'])
                        # 将多级嵌套列表转换成字符串
                        shape.text += PPTUtil.get_list_str('', item['content'], item['children']) + '\n'
                    # 处理图片
                    if "Picture Placeholder" in placeholder:
                        fileName = PPTUtil.get_no_format_or_placeholder(item['content'])
                        image_path = './' + fileName
                        # 处理图片地址包含其它文字的情况
                        if not os.path.exists(image_path):
                            pattern = r'(temp.*)'  # 强制路径以 temp 开头
                            match = re.search(pattern, fileName)
                            if match:
                                fileName = match.group(1)
                                image_path = './' + fileName
                        if os.path.exists(image_path):
                            with Image.open(image_path) as img:
                                # 获取像素尺寸
                                width_px, height_px = img.size

                                # 获取 DPI 信息（如果存在）
                                dpi = img.info.get('dpi', (72, 72))  # 默认 DPI 为 72
                                dpi_x, dpi_y = dpi

                                # 计算英寸尺寸
                                width_inches = width_px / dpi_x
                                height_inches = height_px / dpi_y

                                # 转换为 EMU
                                image_width = Inches(width_inches)
                                image_height = Inches(height_inches)

                            placeholder_left = shape.left
                            placeholder_top = shape.top
                            placeholder_width = shape.width
                            placeholder_height = shape.height

                            aspect_ratio = image_width / image_height

                            if placeholder_width / placeholder_height > aspect_ratio:
                                # 图片相对较高，以高度为基准
                                new_height = placeholder_height
                                new_width = int(new_height * aspect_ratio)
                                placeholder_left = (placeholder_width - new_width) / 2 + placeholder_left
                            else:
                                # 图片相对较宽，以宽度为基准
                                new_width = placeholder_width
                                new_height = int(new_width / aspect_ratio)
                                placeholder_top = (placeholder_height - new_height) / 2 + placeholder_top

                            # 删除原始占位符
                            shape._element.getparent().remove(shape._element)

                            # 插入图片（先使用占位符的尺寸）
                            placeholder_picture = slide.shapes.add_picture(
                                image_path,
                                left=placeholder_left,
                                top=placeholder_top,
                                width=new_width,
                                height=new_height
                            )
    # 提取版式或占位符
    @staticmethod
    def get_format_or_placeholder(text:str):
        index = PPTUtil.get_last_bracket_index(text)
        # 处理提取为空的情况
        if index==None:
            return ""
        return text[index+1:-1]
    # 获取最后一个左括号出现的位置
    @staticmethod
    def get_last_bracket_index(text:str):
        stack = []
        stack.append(text[-1])
        for i in range(len(text) - 2, -1, -1):
            if text[i]=='(' or text[i]=='（':
                stack.pop()
                if len(stack) == 0:
                    return i
            if text[i]==')' or text[i]=='）':
                stack.append(text[i])

    # 获取不带版式或占位符的文本内容
    @staticmethod
    def get_no_format_or_placeholder(text: str):
        index = PPTUtil.get_last_bracket_index(text)
        # 处理提取失败的情况
        if index==None:
            return text
        return text[:index]
    # 保存ppt
    @staticmethod
    def save_ppt(prs,path):
        prs.save(path)
    # 获取模板PPT中的所有版式以及版式的占位符
    @staticmethod
    def get_template_info(template=0):
        # 获取ppt文件，其中template指定使用的模板
        fileName = "./template/Template" + str(template) + ".pptx"
        prs = Presentation(fileName)
        template_info = {}
        for i, layout in enumerate(prs.slide_layouts):
            slide = prs.slides.add_slide(layout)
            placeholder_name = []
            for placeholder in slide.placeholders:
                placeholder_name.append(placeholder.name)
            template_info[layout.name]=placeholder_name
        PPTUtil.template_info = template_info
        return template_info
    # 解析PPT制作方案中需要生成的图片相关信息，并且保存至image_data（于sessionID绑定），返回的生成方案中图片下载地址替换成真实路径
    @staticmethod
    def parse_image(ppt_txt: str,image_data,user_input):
        lines = ppt_txt.splitlines()
        new_ppt_txt = ""
        hava_image = False
        for line in lines:
            new_line = line+'\n'
            # 根据占位符寻找需要生成的图片
            if "Picture Placeholder" in line:
                # 解析文本内容
                content = PPTUtil.get_no_format_or_placeholder(line)
                placeholder = PPTUtil.get_format_or_placeholder(line)
                # 去除-符号
                sub_content = re.sub(r'^-\s*', '', content)
                # 保存-以及后面的空格方便还原最后的PPT制作方案中的图片文字内容
                match = re.match(r'^(-\s*)', content)
                head_content = match.group(1) if match else ''
                # 根据";下载地址:"来分割内容
                split_flag = ";下载地址:"
                if split_flag in sub_content:
                    image_des = sub_content.split(";下载地址:")[0]
                    image_url = sub_content.split(";下载地址:")[1]
                    # 根据"none"判断该图片是否处理过
                    if "none" in image_url:
                        # 生成图片保存路径，并且将图片描述以及保存路径以及相关信息根据sessionId保存
                        image_id = str(uuid.uuid4())
                        image_path = 'temp/'+image_id+'.png'
                        new_line = head_content+image_des+";下载地址:"+image_path+'('+placeholder+')'

                        data ={}
                        data["save_path"]=image_path
                        data["image_description"]=image_des
                        data["context"]=user_input
                        data["is_make"]=False

                        image_data[image_id]=data
                    else:
                        # 设置处理过并且仍然未删除的图片
                        index = image_url.find("temp/")
                        if index!=-1:
                            fileName = image_url[index+len("temp/"):]
                            image_id = os.path.splitext(fileName)[0]
                            if image_id in image_data:
                                image_data[image_id]["is_make"] = False
                    hava_image = True

            new_ppt_txt+=new_line
        return new_ppt_txt,hava_image
    # 判断首次回答的PPT制作方案是否合理
    @staticmethod
    def examine_ppt_scheme_use_ai(ppt_scheme,prompt,service_name):
        llm_input = prompt.format(
            original_scheme=ppt_scheme
        )
        buffer = ""
        is_in_special = False  # 标记此时是否处于特殊数据块
        for chunk in AiforceUtil.chat(llm_input, service_name, str(uuid.uuid4())):
            if not is_in_special:
                # 使用正则表达式查找匹配位置
                match = re.search(r'</think>\s*```\s*ppt', buffer)
                if match:
                    yield chunk
            buffer += chunk
        # 使用正则表达式查找匹配位置
        match = re.search(r'</think>\s*```\s*ppt', buffer)
        if not match:
            # 抛出没有开始符号异常
            raise ValueError("ppt方案生成失败，请重新尝试")






